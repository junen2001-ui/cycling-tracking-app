# -*- coding: utf-8 -*-
"""
要求仕様デジタル化(YAML化)による開発プロセス改革 説明資料
python-pptx で .pptx を生成する
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------- カラーパレット ----------
NAVY      = RGBColor(0x1F, 0x3A, 0x5F)
BLUE      = RGBColor(0x2E, 0x86, 0xAB)
LIGHTBLUE = RGBColor(0xE8, 0xF1, 0xF7)
GRAY      = RGBColor(0x4A, 0x4A, 0x4A)
LIGHTGRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
RED       = RGBColor(0xC0, 0x39, 0x2B)
GREEN     = RGBColor(0x1E, 0x8A, 0x4C)
ORANGE    = RGBColor(0xE0, 0x8E, 0x0B)
DARKGRAY  = RGBColor(0x33, 0x33, 0x33)

FONT_JP = "游ゴシック"
FONT_JP_BOLD = "游ゴシック"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

PAGE_NUM = {"n": 0}

# ---------- 共通ユーティリティ ----------

def add_slide():
    s = prs.slides.add_slide(BLANK)
    # 背景
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s

def set_font(run, size=18, color=DARKGRAY, bold=False, font=FONT_JP, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)

def add_textbox(slide, left, top, width, height, text, size=18, color=DARKGRAY,
                 bold=False, align=PP_ALIGN.LEFT, font=FONT_JP, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold, font, italic)
    return tb

def add_footer(slide, page_number=None):
    if page_number is None:
        PAGE_NUM["n"] += 1
        n = PAGE_NUM["n"]
    else:
        n = page_number
    # 下部バー
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.28), SLIDE_W, Inches(0.28))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "要求仕様デジタル化プロジェクト"
    set_font(run, 9, WHITE, False)
    # ページ番号
    pnbox = slide.shapes.add_textbox(SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.28), Inches(0.75), Inches(0.28))
    tf2 = pnbox.text_frame
    tf2.margin_top = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(n)
    set_font(run2, 10, WHITE, True)

def add_title_bar(slide, title, subtitle=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_font(run, 26, WHITE, True)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(1.05), SLIDE_W - Inches(0.8), Inches(0.4),
                    subtitle, size=14, color=BLUE, bold=True)
        return Inches(1.55)
    return Inches(1.25)

def add_bullets(slide, left, top, width, height, items, size=16, color=DARKGRAY, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
        prefix = "● " if level == 0 else "- "
        indent = Inches(0.0) if level == 0 else Inches(0.35)
        p.text = ""
        run = p.add_run()
        run.text = prefix + text
        sz = size if level == 0 else size - 2
        col = color if level == 0 else GRAY
        set_font(run, sz, col, bold=(level == 0))
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(indent))
    return tb

def add_box(slide, left, top, width, height, text, fill=LIGHTBLUE, text_color=NAVY,
            size=14, bold=True, line_color=BLUE, align=PP_ALIGN.CENTER, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line_color
    box.line.width = Pt(1.25)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, text_color, bold)
    return box

def add_arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn

def add_code_box(slide, left, top, width, height, lines, size=13, title=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x27, 0x2B, 0x30)
    box.line.color.rgb = RGBColor(0x11, 0x13, 0x15)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.1)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = 1.05
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        set_font(run, size, RGBColor(0xB5, 0xE8, 0x53), False, font="Consolas")
    return box

def add_table(slide, left, top, width, height, data, col_widths=None, header_fill=NAVY,
              header_color=WHITE, body_size=13, header_size=13, row_h=None):
    rows = len(data)
    cols = len(data[0])
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = w
    if row_h:
        for r in range(rows):
            gtable.rows[r].height = row_h
    for r in range(rows):
        for c in range(cols):
            cell = gtable.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            if r == 0:
                set_font(run, header_size, header_color, True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                set_font(run, body_size, DARKGRAY, False)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHTGRAY
    return gtable

# ============================================================
# Slide 1: タイトル
# ============================================================
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.line.fill.background()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.55), SLIDE_W, Inches(0.06))
accent.fill.solid()
accent.fill.fore_color.rgb = BLUE
accent.line.fill.background()
add_textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0),
            "要求仕様デジタル化による開発プロセス改革", size=34, color=WHITE, bold=True)
add_textbox(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.7),
            "ExcelからYAMLへ ― AIが理解できる要求仕様が拓く、品質と開発効率の刷新",
            size=18, color=RGBColor(0xB8, 0xD4, 0xEA), bold=False)
add_textbox(s, Inches(0.9), Inches(6.6), Inches(6), Inches(0.4),
            "組み込みファームウェア開発チーム", size=13, color=RGBColor(0x9A, 0xB4, 0xC8))

# ============================================================
# Slide 2: Why-1 現状の開発プロセス
# ============================================================
s = add_slide()
top = add_title_bar(s, "現状の開発プロセス", "ウォーターフォール型 + 度重なる仕様変更")
steps = ["要求仕様\n(Excel/Word)", "基本設計", "機能設計", "Coding", "単体テスト\n(CT)", "組込みテスト\n(IT)", "システム\nテスト(ST)"]
n = len(steps)
box_w = Inches(1.55)
gap = Inches(0.15)
total_w = box_w * n + gap * (n - 1)
start_x = (SLIDE_W - total_w) // 2
y = Inches(2.3)
xs = []
for i, st in enumerate(steps):
    x = start_x + i * (box_w + gap)
    xs.append(x)
    add_box(s, x, y, box_w, Inches(0.9), st, fill=LIGHTBLUE, text_color=NAVY, size=12)
    if i < n - 1:
        add_arrow(s, x + box_w, y + Inches(0.45), x + box_w + gap, y + Inches(0.45))
# 仕様変更の差し戻し矢印
add_arrow(s, xs[4] + Inches(0.6), y + Inches(1.0), xs[1] + Inches(0.8), y + Inches(1.0), color=RED, width=2.0)
rework = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xs[2], y + Inches(0.95), Inches(2.3), Inches(0.02))
add_textbox(s, xs[1], y + Inches(1.05), Inches(3.5), Inches(0.4), "仕様変更 → 設計手戻り(繰り返し発生)", size=13, color=RED, bold=True)
add_bullets(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(2.8), [
    "基本的にはウォーターフォールモデルだが、開発途中で仕様変更が発生し、設計変更を何度も繰り返している",
    "上流の「要求仕様」が Excel / Word というアナログ媒体であることが、この繰り返しのたびに問題を増幅させている",
], size=17)

# ============================================================
# Slide 3: Why-2 課題1・2
# ============================================================
s = add_slide()
add_title_bar(s, "現状の3つの課題(1/2)", "アナログな要求仕様が引き起こす問題")
add_box(s, Inches(0.9), Inches(1.9), Inches(0.6), Inches(0.6), "1", fill=RED, text_color=WHITE, size=22, shape=MSO_SHAPE.OVAL)
add_textbox(s, Inches(1.7), Inches(1.85), Inches(10.5), Inches(0.5), "変更点の見落とし", size=20, color=NAVY, bold=True)
add_bullets(s, Inches(1.7), Inches(2.4), Inches(10.5), Inches(1.3), [
    "変更点はExcel/Wordを目視で比較するしかなく、確認は担当者の注意力に依存する",
    "セルの色付けやコメントに頼った運用では、更新漏れ・見落としが構造的に発生する",
], size=15)
add_box(s, Inches(0.9), Inches(4.0), Inches(0.6), Inches(0.6), "2", fill=RED, text_color=WHITE, size=22, shape=MSO_SHAPE.OVAL)
add_textbox(s, Inches(1.7), Inches(3.95), Inches(10.5), Inches(0.5), "影響範囲の見極めが困難", size=20, color=NAVY, bold=True)
add_bullets(s, Inches(1.7), Inches(4.5), Inches(10.5), Inches(1.6), [
    "要求と設計・テストの対応関係が文書として構造化されていないため、変更がどこに波及するか人手で追うしかない",
    "結果として、リグレッションテストで確認すべき項目が漏れ、デグレードが本番間際・出荷後に発覚するリスクが残る",
], size=15)

# ============================================================
# Slide 4: Why-3 課題3
# ============================================================
s = add_slide()
add_title_bar(s, "現状の3つの課題(2/2)", "小手先修正の連鎖がソースを複雑化させる")
add_box(s, Inches(0.9), Inches(1.9), Inches(0.6), Inches(0.6), "3", fill=RED, text_color=WHITE, size=22, shape=MSO_SHAPE.OVAL)
add_textbox(s, Inches(1.7), Inches(1.85), Inches(10.5), Inches(0.5), "納期優先の小手先修正が繰り返される", size=20, color=NAVY, bold=True)
add_bullets(s, Inches(1.7), Inches(2.4), Inches(10.5), Inches(1.3), [
    "納期が迫ると、根本原因に手を入れる時間が取れず、局所的なパッチで対応してしまう",
    "この対応が積み重なることで、ソースコードの構造が徐々に複雑化し、次の変更コストがさらに上がる",
], size=15)
# 悪循環図
cx, cy, r = Inches(8.6), Inches(5.55), Inches(1.35)
labels = ["納期のプレッシャー", "小手先の修正", "コードの複雑化", "変更コストの増大"]
positions = [(0,-1),(1,0),(0,1),(-1,0)]
box_positions = []
for i,(dx,dy) in enumerate(positions):
    bx = cx + Emu(int(dx * r))
    by = cy + Emu(int(dy * r * 0.85))
    box_positions.append((bx,by))
bw, bh = Inches(2.1), Inches(0.6)
for i,(bx,by) in enumerate(box_positions):
    add_box(s, bx - bw//2, by - bh//2, bw, bh, labels[i], fill=RGBColor(0xFC,0xE8,0xE6), text_color=RED, size=12)
edge_offset = 950000  # EMU, ボックス外周から矢印を出すためのオフセット
for i in range(4):
    x1,y1 = box_positions[i]
    x2,y2 = box_positions[(i+1)%4]
    dx, dy = int(x2) - int(x1), int(y2) - int(y1)
    dist = math.hypot(dx, dy)
    ux, uy = dx / dist, dy / dist
    sx = Emu(int(int(x1) + ux * edge_offset))
    sy = Emu(int(int(y1) + uy * edge_offset))
    ex = Emu(int(int(x2) - ux * edge_offset))
    ey = Emu(int(int(y2) - uy * edge_offset))
    add_arrow(s, sx, sy, ex, ey, color=RED, width=1.75)
add_textbox(s, Inches(0.9), Inches(4.9), Inches(4.7), Inches(1.8),
            "この3つの課題はいずれも、要求仕様が「機械可読でない」ことに起因している。\n\n"
            "→ 上流の要求仕様そのものをデジタル化し、変更・影響範囲・品質を機械的に扱える形にすることが、根本解決の出発点となる。",
            size=14, color=DARKGRAY)

for sl in prs.slides.__iter__():
    pass

# ============================================================
# Slide 5: What-1 目指す全体像
# ============================================================
s = add_slide()
add_title_bar(s, "目指す姿:要求仕様を起点とした自動化パイプライン")
flow = ["要求仕様\n(YAML)", "Git\n差分検出", "設計項目\n自動生成", "CT/IT\n項目自動生成", "インシデント\n管理へ自動登録"]
n = len(flow)
box_w = Inches(2.1); gap = Inches(0.28)
total_w = box_w*n + gap*(n-1)
start_x = (SLIDE_W - total_w)//2
y = Inches(2.5)
for i, st in enumerate(flow):
    x = start_x + i*(box_w+gap)
    fill = LIGHTBLUE if i>0 else RGBColor(0xD9,0xEA,0xD3)
    tc = NAVY if i>0 else GREEN
    add_box(s, x, y, box_w, Inches(1.0), st, fill=fill, text_color=tc, size=13)
    if i < n-1:
        add_arrow(s, x+box_w, y+Inches(0.5), x+box_w+gap, y+Inches(0.5), width=2.25)
add_textbox(s, start_x, y+Inches(1.25), total_w, Inches(0.4), "版数管理・差分・自動生成をすべて機械的に処理", size=13, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.5), [
    "要求仕様をYAML(構造化テキスト)で記述し、Gitでバージョン管理する",
    "変更はGit差分として機械的に検出でき、影響範囲の特定を人手に頼らない",
    "差分から設計項目・CT/ITのテストケースを自動生成し、レビューとテスト双方の抜け漏れを防ぐ",
    "検出した課題や不整合はインシデント管理(バックログ等)へ自動登録し、対応漏れを防止する",
], size=15)

# ============================================================
# Slide 6: What-2 Before/After比較
# ============================================================
s = add_slide()
add_title_bar(s, "Before / After", "何が変わり、何が変わらないか")
data = [
    ["項目", "Before(現状)", "After(YAML化後)"],
    ["要求仕様の記述", "Excel / Word", "YAML(構造化テキスト)"],
    ["変更点の検出", "目視での比較", "Git差分による機械的検出"],
    ["影響範囲の特定", "人手による調査", "要求ID紐付けによる自動トレース"],
    ["設計・テスト項目", "都度手作業で作成", "差分から自動生成"],
    ["版数管理", "ファイル名・フォルダ運用", "Gitによる正式な版数管理"],
    ["設計の考え方・レビュー観点", "変わらない", "変わらない"],
]
add_table(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.6), data,
          col_widths=[Inches(3.3), Inches(4.1), Inches(4.1)], row_h=Inches(0.62))
add_textbox(s, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.5),
            "変わるのは「記述の形式」。設計スキルやレビューの目的は今まで通り活きる。", size=14, color=NAVY, bold=True)

# ============================================================
# Slide 7: YAMLとは-1 基本
# ============================================================
s = add_slide()
add_title_bar(s, "YAMLってどんなもの?", "テキストで書ける、人にもAIにも読みやすいデータ形式")
add_bullets(s, Inches(0.9), Inches(1.9), Inches(5.7), Inches(4.5), [
    "YAML(YAML Ain't Markup Language)は、構造化データをテキストで表現するための記述形式",
    "「キー: 値」の組でデータを表し、字下げ(インデント)で階層構造を表現する",
    "リスト(繰り返し項目)は行頭の「-」で表す",
    "タブ文字は使えず、同じ階層は必ず同じ数の半角スペースで揃える",
    "拡張子は .yaml または .yml",
    "Docker Compose・Kubernetes・GitHub Actionsなど、多くの開発ツールの設定ファイルとして既に広く使われている",
], size=15)
add_code_box(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(3.0), [
    "# センサー設定の例",
    "sensor:",
    "  name: TempSensor01",
    "  type: I2C",
    "  address: 0x48",
    "  sampling_rate_ms: 100",
    "  channels:",
    "    - id: 0",
    "      unit: degC",
    "    - id: 1",
    "      unit: percentRH",
], size=14)
add_textbox(s, Inches(6.9), Inches(5.05), Inches(5.5), Inches(0.5),
            "← インデントと「-」だけで、階層とリストを表現できる", size=12, color=GRAY)

# ============================================================
# Slide 8: YAMLとは-2 Excelとの対比
# ============================================================
s = add_slide()
add_title_bar(s, "Excelの表と何が違うのか", "見た目は違うが、伝えている情報は同じ")
add_box(s, Inches(0.9), Inches(1.9), Inches(5.5), Inches(0.5), "Excel(現状)", fill=RGBColor(0xF6,0xE3,0xC6), text_color=ORANGE, size=15)
data2 = [["項番","状態","イベント","遷移先"],["1","IDLE","Start押下","RUNNING"],["2","RUNNING","異常検知","ERROR"]]
add_table(s, Inches(0.9), Inches(2.5), Inches(5.5), Inches(1.8), data2,
          col_widths=[Inches(0.9),Inches(1.5),Inches(1.6),Inches(1.5)], row_h=Inches(0.5))
add_textbox(s, Inches(0.9), Inches(4.5), Inches(5.5), Inches(1.6),
            "見た目は整っているが、セル結合・自由記述・表記ゆれによりプログラムからは正確に読み取れない", size=13, color=GRAY)
add_box(s, Inches(6.9), Inches(1.9), Inches(5.5), Inches(0.5), "YAML(After)", fill=LIGHTBLUE, text_color=NAVY, size=15)
add_code_box(s, Inches(6.9), Inches(2.5), Inches(5.5), Inches(2.2), [
    "state_machine:",
    "  - state: IDLE",
    "    on: {event: START, next: RUNNING}",
    "  - state: RUNNING",
    "    on: {event: ERROR_DETECTED, next: ERROR}",
], size=14)
add_textbox(s, Inches(6.9), Inches(4.9), Inches(5.5), Inches(1.4),
            "同じ内容を、キーと値の組で厳密に表現。表記ゆれが起きず、プログラム(AI含む)が確実に解析できる",
            size=13, color=GRAY)

# ============================================================
# Slide 9: YAMLとは-3 参考URL
# ============================================================
s = add_slide()
add_title_bar(s, "まずは触ってみる:参考リソース", "初めてでも読める入門コンテンツ")
add_textbox(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(0.4), "YAML入門(初心者向け)", size=17, color=NAVY, bold=True)
add_bullets(s, Inches(1.0), Inches(2.35), Inches(11.3), Inches(1.9), [
    "【初心者】YAML入門 (Qiita) - qiita.com/mackeyTA/items/7dd9282ae7c0599495a6",
    "YAML入門:設定ファイルでよく見る「読みやすい」書き方 (Zenn) - zenn.dev/acntechjp/articles/a39e710de5d744",
    "YAML入門:サンプルから学べる初心者向けガイド (CircleCI公式ブログ) - circleci.com/blog/what-is-yaml-a-beginner-s-guide/",
], size=14)
add_textbox(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.4), "組み込み分野での実例", size=17, color=NAVY, bold=True)
add_bullets(s, Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.3), [
    "Devicetree versus Kconfig (Zephyr RTOS 公式ドキュメント) - docs.zephyrproject.org/latest/build/dts/dt-vs-kconfig.html",
    "  組み込みOSであるZephyrは、ハードウェア記述(devicetree)をYAML bindingsで管理しており、実際の量産組み込み開発でのYAML活用例として参考になる",
], size=14)
add_textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.7),
            "まずは既存のYAMLファイル(CI設定など)を眺めてみるだけでも感覚がつかめる。難しい構文を覚える前に「読む」ことから始めればよい。",
            size=13, color=BLUE, bold=True)

# ============================================================
# Slide 10: 記述手法の全体マップ
# ============================================================
s = add_slide()
add_title_bar(s, "要求記述手法の全体マップ", "目的に応じて9つの手法を使い分ける")
data3 = [
    ["手法", "主な目的"],
    ["判断表(デシジョンテーブル)", "条件の組み合わせと動作を漏れなく整理"],
    ["EARS", "自然文の要求を曖昧さのない定型構文で記述"],
    ["正規化テーブル", "パラメータ・データ項目を単一の真実源として管理"],
    ["シーケンス表", "処理・メッセージのやり取りの順序を記述"],
    ["状態遷移表", "状態とイベントによる振る舞いの変化を記述"],
    ["タイミングチャート", "割込み・周期処理などハードウェアのタイミング関係を記述"],
    ["I/Oマップ/レジスタ定義表", "ピン配置・レジスタビット割付を管理し、コード生成につなげる"],
    ["Given-When-Then", "テスト観点で要求を記述し、CT/ITケースへ橋渡しする"],
    ["トレーサビリティマトリクス", "要求⇔設計⇔テストの対応関係を管理し、影響範囲を追跡"],
]
add_table(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(5.15), data3,
          col_widths=[Inches(3.8), Inches(7.7)], row_h=Inches(0.5), body_size=13)

for sl in prs.slides.__iter__():
    pass

def method_slide(no_total, title, subtitle, bullets, code_lines, code_title=None):
    s = add_slide()
    add_title_bar(s, title, subtitle)
    add_bullets(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(4.8), bullets, size=14.5)
    if code_title:
        add_textbox(s, Inches(6.75), Inches(1.85), Inches(5.7), Inches(0.35), code_title, size=13, color=BLUE, bold=True)
        code_top = Inches(2.2)
    else:
        code_top = Inches(1.95)
    add_code_box(s, Inches(6.75), code_top, Inches(5.7), Inches(4.6), code_lines, size=12.5)
    return s

# Slide 11: 判断表
method_slide(9, "① 判断表(デシジョンテーブル)", "条件の組み合わせと動作を漏れなく整理する",
    [
        "複数の条件が絡む判定ロジックを、条件×動作のマトリクスで網羅的に整理する手法",
        "「組み合わせの漏れ」を防げるため、複雑な起動条件・異常判定ロジックとの相性が良い",
        "YAMLでは conditions と rules のリストとして表現し、条件の組み合わせパターンを列挙する",
    ],
    [
        "decision_table:",
        "  name: モータ起動判定",
        "  conditions:",
        "    - door_closed",
        "    - battery_ok",
        "    - start_button",
        "  rules:",
        "    - when: [true, true, true]",
        "      action: MOTOR_START",
        "    - when: [false, true, true]",
        "      action: MOTOR_STOP",
        "    - when: [true, false, true]",
        "      action: MOTOR_STOP",
    ], "YAML記述例")

# Slide 12: EARS
method_slide(9, "② EARS", "曖昧さのない定型構文で要求を書く",
    [
        "Easy Approach to Requirements Syntax の略。自然文の要求を定型パターンに当てはめて記述する手法",
        "「WHEN〜SHALL」「IF〜THEN SHALL」のように主語・条件・義務を明確化し、解釈のブレをなくす",
        "Event-driven / State-driven / Unwanted behavior 等のパターンをYAMLの pattern 属性で明示する",
    ],
    [
        "requirements:",
        "  - id: REQ-101",
        "    pattern: event-driven",
        "    text: >",
        "      WHEN the start button is",
        "      pressed, the system SHALL",
        "      start the motor within 100ms",
        "  - id: REQ-102",
        "    pattern: unwanted-behavior",
        "    text: >",
        "      IF battery voltage < 3.0V,",
        "      THEN the system SHALL enter",
        "      safe mode",
    ], "YAML記述例")

# Slide 13: 正規化テーブル
method_slide(9, "③ 正規化テーブル", "パラメータを単一の真実源(Single Source of Truth)で管理",
    [
        "定数・パラメータ・データ項目を一箇所で一意に定義し、重複や表記ゆれを排除する",
        "同じ値がExcelの複数シートに散在する状態を解消し、変更時の修正漏れを防ぐ",
        "設計書・コード双方から同じ定義を参照できるため、値のズレによる不具合を根本から防止できる",
    ],
    [
        "parameters:",
        "  - id: P001",
        "    name: MAX_MOTOR_SPEED",
        "    unit: rpm",
        "    value: 3000",
        "    range: [0, 5000]",
        "  - id: P002",
        "    name: SAMPLING_RATE_HZ",
        "    unit: Hz",
        "    value: 100",
        "    range: [10, 1000]",
    ], "YAML記述例")

# Slide 14: シーケンス表
method_slide(9, "④ シーケンス表", "処理・メッセージのやり取りの順序を記述する",
    [
        "モジュール間・タスク間でやり取りされる処理やメッセージの順序を時系列に整理する",
        "起動シーケンスや異常時のフェイルセーフ処理など、順序が重要な処理の記述に適する",
        "YAMLでは steps のリストとして、from / to / message を順に並べて表現する",
    ],
    [
        "sequence:",
        "  name: 起動シーケンス",
        "  steps:",
        "    - step: 1",
        "      from: MainTask",
        "      to: MotorDriver",
        "      message: INIT",
        "    - step: 2",
        "      from: MotorDriver",
        "      to: MainTask",
        "      message: INIT_DONE",
        "    - step: 3",
        "      from: MainTask",
        "      to: MotorDriver",
        "      message: START",
    ], "YAML記述例")

# Slide 15: 状態遷移表
method_slide(9, "⑤ 状態遷移表", "状態とイベントによる振る舞いの変化を記述する",
    [
        "組み込み制御の中核となる「状態×イベント→遷移先」の対応関係を網羅的に整理する",
        "モータ制御・通信プロトコルなど、状態を持つ制御ロジックのほぼすべてに適用できる",
        "YAMLでは states と transitions のリストで、from / event / to を明示する",
    ],
    [
        "state_machine:",
        "  name: MotorControl",
        "  states: [IDLE, RUNNING, ERROR]",
        "  transitions:",
        "    - from: IDLE",
        "      event: START",
        "      to: RUNNING",
        "    - from: RUNNING",
        "      event: FAULT_DETECTED",
        "      to: ERROR",
        "    - from: ERROR",
        "      event: RESET",
        "      to: IDLE",
    ], "YAML記述例")

# Slide 16: タイミングチャート
method_slide(9, "⑥ タイミングチャート", "割込み・周期処理などハードウェアのタイミングを記述する",
    [
        "組み込み特有の要求記述。信号の変化タイミングや周期処理の時間関係を数値で厳密に表す",
        "状態遷移表だけでは表現しきれない「いつ・どれだけの時間で」という制約を補完する",
        "YAMLでは signal ごとに period や events(時刻とレベル)のリストとして記述する",
    ],
    [
        "timing:",
        "  signal: PWM_OUT",
        "  period_us: 1000",
        "  events:",
        "    - t_us: 0",
        "      level: HIGH",
        "    - t_us: 300",
        "      level: LOW",
        "  constraints:",
        "    - name: min_pulse_width",
        "      value_us: 50",
    ], "YAML記述例")

# Slide 17: I/Oマップ・レジスタ定義表
method_slide(9, "⑦ I/Oマップ / レジスタ定義表", "ピン配置・レジスタビット割付を管理する",
    [
        "組み込み特有の要求記述。ピン配置やレジスタのビット割付をYAMLで一元管理する",
        "この定義から、Cヘッダファイルやconfigコードを自動生成でき、手打ち転記によるミスを撲滅できる",
        "ハードウェア変更時も、この1ファイルを直せば設計・コードへ機械的に反映できる",
    ],
    [
        "io_map:",
        "  pins:",
        "    - pin: PA0",
        "      function: PWM_OUT",
        "      direction: OUTPUT",
        "  registers:",
        "    - name: TIMER1_CTRL",
        "      address: 0x4000_0000",
        "      bits:",
        "        - bit: 3",
        "          name: TIMER1_EN",
        "          access: RW",
    ], "YAML記述例")

# Slide 18: Given-When-Then
method_slide(9, "⑧ Given-When-Then", "テスト観点で要求を記述し、CT/ITへ橋渡しする",
    [
        "振る舞い駆動開発(BDD)で使われる記法。前提(Given)・操作(When)・結果(Then)で要求を記述する",
        "EARSで書いた要求を、そのままテストシナリオの形に変換しやすいのが特長",
        "この構造は単体テスト(CT)・組込みテスト(IT)のテストケース自動生成の入力として直接利用できる",
    ],
    [
        "scenario:",
        "  id: TC-201",
        "  req_id: REQ-101",
        "  given: >",
        "    モータが停止している",
        "  when: >",
        "    Startボタンが押される",
        "  then: >",
        "    100ms以内にモータが",
        "    回転を始める",
    ], "YAML記述例")

# Slide 19: トレーサビリティマトリクス
method_slide(9, "⑨ トレーサビリティマトリクス", "要求⇔設計⇔テストの対応関係を管理する",
    [
        "各要求がどの設計・どのテストケースに対応するかを機械可読な形で紐付ける",
        "「影響範囲の見極めが難しい」という現状課題(課題②)に最も直接的に効く手法",
        "要求が変更されると、紐付いた設計・テストIDを自動的にリストアップでき、対応漏れを防止できる",
    ],
    [
        "traceability:",
        "  - req_id: REQ-101",
        "    design_id: DES-045",
        "    code_ref: motor_ctrl.c#L120",
        "    test_ids:",
        "      - CT-301",
        "      - IT-118",
        "  - req_id: REQ-102",
        "    design_id: DES-046",
        "    test_ids:",
        "      - CT-302",
    ], "YAML記述例")

# ============================================================
# Slide 20: 逆変換デモ-1 状態遷移図
# ============================================================
s = add_slide()
add_title_bar(s, "レビュー時の逆変換デモ(1)", "YAML → 状態遷移図を自動生成")
add_textbox(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(0.35), "書くのは:YAML", size=13, color=BLUE, bold=True)
add_code_box(s, Inches(0.9), Inches(2.2), Inches(5.6), Inches(3.4), [
    "state_machine:",
    "  name: MotorControl",
    "  states: [IDLE, RUNNING, ERROR]",
    "  transitions:",
    "    - from: IDLE",
    "      event: START",
    "      to: RUNNING",
    "    - from: RUNNING",
    "      event: FAULT_DETECTED",
    "      to: ERROR",
    "    - from: ERROR",
    "      event: RESET",
    "      to: IDLE",
], size=13)
add_textbox(s, Inches(6.75), Inches(1.85), Inches(5.7), Inches(0.35), "レビュー時に見るのは:自動生成された状態遷移図", size=13, color=GREEN, bold=True)
diag = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), Inches(2.2), Inches(5.7), Inches(4.3))
diag.fill.solid(); diag.fill.fore_color.rgb = LIGHTGRAY; diag.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
idle = add_box(s, Inches(7.1), Inches(2.55), Inches(1.8), Inches(0.8), "IDLE", fill=LIGHTBLUE, text_color=NAVY, size=15, shape=MSO_SHAPE.OVAL)
running = add_box(s, Inches(10.1), Inches(2.55), Inches(1.9), Inches(0.8), "RUNNING", fill=RGBColor(0xD9,0xEA,0xD3), text_color=GREEN, size=15, shape=MSO_SHAPE.OVAL)
error = add_box(s, Inches(8.6), Inches(4.7), Inches(1.9), Inches(0.8), "ERROR", fill=RGBColor(0xFC,0xE0,0xDD), text_color=RED, size=15, shape=MSO_SHAPE.OVAL)
add_arrow(s, Inches(8.9), Inches(2.95), Inches(10.1), Inches(2.95), color=NAVY, width=2)
add_textbox(s, Inches(8.85), Inches(2.55), Inches(1.3), Inches(0.3), "START", size=10, color=NAVY, align=PP_ALIGN.CENTER)
add_arrow(s, Inches(10.6), Inches(3.35), Inches(9.6), Inches(4.7), color=NAVY, width=2)
add_textbox(s, Inches(9.7), Inches(3.9), Inches(1.8), Inches(0.5), "FAULT_DETECTED", size=10, color=NAVY, align=PP_ALIGN.CENTER)
add_arrow(s, Inches(8.6), Inches(5.0), Inches(7.6), Inches(3.35), color=NAVY, width=2)
add_textbox(s, Inches(7.2), Inches(4.0), Inches(1.3), Inches(0.3), "RESET", size=10, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.9),
            "記述はYAML、確認は従来通りの状態遷移図。レビュー担当者はツールの出す図を見るだけでよく、YAMLの読み書きスキルを必須としない。",
            size=13, color=GRAY)

# ============================================================
# Slide 21: 逆変換デモ-2 判断表
# ============================================================
s = add_slide()
add_title_bar(s, "レビュー時の逆変換デモ(2)", "YAML → 判断表(表形式)を自動生成")
add_textbox(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(0.35), "書くのは:YAML", size=13, color=BLUE, bold=True)
add_code_box(s, Inches(0.9), Inches(2.2), Inches(5.6), Inches(3.4), [
    "decision_table:",
    "  name: モータ起動判定",
    "  conditions:",
    "    - door_closed",
    "    - battery_ok",
    "    - start_button",
    "  rules:",
    "    - when: [true, true, true]",
    "      action: MOTOR_START",
    "    - when: [false, true, true]",
    "      action: MOTOR_STOP",
    "    - when: [true, false, true]",
    "      action: MOTOR_STOP",
], size=13)
add_textbox(s, Inches(6.75), Inches(1.85), Inches(5.7), Inches(0.35), "レビュー時に見るのは:自動生成された判断表", size=13, color=GREEN, bold=True)
data4 = [
    ["door_closed","battery_ok","start_button","action"],
    ["TRUE","TRUE","TRUE","MOTOR_START"],
    ["FALSE","TRUE","TRUE","MOTOR_STOP"],
    ["TRUE","FALSE","TRUE","MOTOR_STOP"],
]
add_table(s, Inches(6.75), Inches(2.25), Inches(5.7), Inches(2.0), data4,
          col_widths=[Inches(1.5),Inches(1.5),Inches(1.5),Inches(1.2)], row_h=Inches(0.5), body_size=12, header_size=12)
add_textbox(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.9),
            "見慣れた判断表の形式で確認できるため、これまでExcelでレビューしてきたメンバーもそのまま移行できる。",
            size=13, color=GRAY)

# ============================================================
# Slide 22: 逆変換デモ-3 タイミングチャート
# ============================================================
s = add_slide()
add_title_bar(s, "レビュー時の逆変換デモ(3)", "YAML → タイミングチャートを自動生成")
add_textbox(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(0.35), "書くのは:YAML", size=13, color=BLUE, bold=True)
add_code_box(s, Inches(0.9), Inches(2.2), Inches(5.6), Inches(3.4), [
    "timing:",
    "  signal: PWM_OUT",
    "  period_us: 1000",
    "  events:",
    "    - t_us: 0",
    "      level: HIGH",
    "    - t_us: 300",
    "      level: LOW",
    "  constraints:",
    "    - name: min_pulse_width",
    "      value_us: 50",
], size=13)
add_textbox(s, Inches(6.75), Inches(1.85), Inches(5.7), Inches(0.35), "レビュー時に見るのは:自動生成された波形図", size=13, color=GREEN, bold=True)
chart_left = Inches(6.9); chart_top = Inches(2.6); chart_h = Inches(1.0)
axis = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, chart_left, chart_top + chart_h, chart_left + Inches(5.3), chart_top + chart_h)
axis.line.color.rgb = DARKGRAY; axis.line.width = Pt(1.25)
high1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top, Inches(1.6), Inches(0.35))
high1.fill.solid(); high1.fill.fore_color.rgb = BLUE; high1.line.fill.background()
low1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, chart_left+Inches(1.6), chart_top+Inches(0.35), chart_left+Inches(1.6), chart_top+chart_h)
low1.line.color.rgb = DARKGRAY; low1.line.width = Pt(1.25)
flat1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, chart_left+Inches(1.6), chart_top+chart_h, chart_left+Inches(5.3), chart_top+chart_h)
flat1.line.color.rgb = DARKGRAY; flat1.line.width = Pt(1.25)
add_textbox(s, chart_left, chart_top+Inches(1.15), Inches(1.8), Inches(0.3), "300us (HIGH)", size=11, color=BLUE)
add_textbox(s, chart_left+Inches(2.0), chart_top+Inches(1.15), Inches(2.6), Inches(0.3), "700us (LOW)", size=11, color=GRAY)
add_textbox(s, Inches(6.9), Inches(4.2), Inches(5.5), Inches(0.4), "period_us = 1000 (周期1ms)", size=12, color=NAVY, bold=True)
add_textbox(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.9),
            "ハードウェアタイミングの制約も波形として可視化され、口頭説明に頼らず正確にレビューできる。",
            size=13, color=GRAY)

for sl in prs.slides.__iter__():
    pass

# ============================================================
# Slide 23: メリット整理-1
# ============================================================
s = add_slide()
add_title_bar(s, "デジタル化がもたらすメリット(1/2)", "版数管理・差分・自動生成の基本メリット")
merits1 = [
    ("Gitによる正しい版数管理", "誰が・いつ・何を変更したかが履歴として残り、レビューと復元が確実になる"),
    ("正確な要求の差分抽出", "Git diffにより変更点を機械的に特定でき、目視確認の見落としを排除する"),
    ("設計・テスト項目の自動生成", "差分から影響を受ける設計項目・テストケースを自動リストアップできる"),
    ("インシデント管理への自動登録", "検出した不整合や要対応項目をバックログ等へ自動起票し、対応漏れを防ぐ"),
]
y = Inches(1.95)
for title, desc in merits1:
    add_box(s, Inches(0.9), y, Inches(0.5), Inches(0.9), "✓", fill=GREEN, text_color=WHITE, size=20, shape=MSO_SHAPE.OVAL)
    add_textbox(s, Inches(1.6), y, Inches(10.6), Inches(0.4), title, size=17, color=NAVY, bold=True)
    add_textbox(s, Inches(1.6), y+Inches(0.42), Inches(10.6), Inches(0.55), desc, size=13, color=GRAY)
    y += Inches(1.15)

# ============================================================
# Slide 24: メリット整理-2
# ============================================================
s = add_slide()
add_title_bar(s, "デジタル化がもたらすメリット(2/2)", "組み込み開発ならではの追加メリット")
merits2 = [
    ("コードの自動生成", "I/Oマップ・状態遷移表からCヘッダやconfigコードを自動生成し、手打ち転記のミスを撲滅する"),
    ("スキーマによる妥当性検証", "JSON Schema等で必須項目漏れ・型不整合を機械的にチェックし、不備を作成時点で検出する"),
    ("AIによる矛盾・抜け漏れ検出", "要求間の矛盾や記載漏れをAIが機械的に検出し、レビュー負荷を下げつつ精度を上げる"),
    ("CI連携による一気通貫の自動化", "要求変更 → 影響設計/テストのリストアップ → レビュー依頼までを自動パイプライン化できる"),
    ("製品ファミリー間での差分再利用", "派生機種開発において、共通要求と差分要求を明確に分離・再利用できる"),
]
y = Inches(1.75)
for title, desc in merits2:
    add_box(s, Inches(0.9), y, Inches(0.5), Inches(0.78), "✓", fill=BLUE, text_color=WHITE, size=18, shape=MSO_SHAPE.OVAL)
    add_textbox(s, Inches(1.6), y-Inches(0.02), Inches(10.6), Inches(0.4), title, size=16, color=NAVY, bold=True)
    add_textbox(s, Inches(1.6), y+Inches(0.38), Inches(10.6), Inches(0.5), desc, size=12.5, color=GRAY)
    y += Inches(1.0)

# ============================================================
# Slide 25: FAQ-1
# ============================================================
def faq_slide(title, subtitle, qa_list):
    s = add_slide()
    add_title_bar(s, title, subtitle)
    y = Inches(1.95)
    for q, a in qa_list:
        add_box(s, Inches(0.9), y, Inches(0.55), Inches(0.55), "Q", fill=ORANGE, text_color=WHITE, size=18, shape=MSO_SHAPE.OVAL)
        add_textbox(s, Inches(1.6), y+Inches(0.02), Inches(10.6), Inches(0.55), q, size=16, color=NAVY, bold=True)
        add_box(s, Inches(0.9), y+Inches(0.75), Inches(0.55), Inches(0.55), "A", fill=BLUE, text_color=WHITE, size=18, shape=MSO_SHAPE.OVAL)
        add_textbox(s, Inches(1.6), y+Inches(0.77), Inches(10.6), Inches(0.75), a, size=14, color=GRAY)
        y += Inches(1.7)
    return s

faq_slide("よくある不安への回答(1/3)", "「今までのやり方が否定される」わけではない", [
    ("今まで培ってきた設計スキルは無駄になるのでは?", "ならない。判断表・状態遷移・シーケンスといった設計の考え方そのものは変わらない。変わるのは、それを書き留める記法だけ。"),
    ("YAMLを書くのはプログラミングみたいで難しいのでは?", "プログラミング言語ではない。「項目名: 値」を字下げして並べるだけで、Excelの表を作る感覚に近い。テンプレートも用意する。"),
])

# ============================================================
# Slide 26: FAQ-2
# ============================================================
faq_slide("よくある不安への回答(2/3)", "レビューの進め方は変わらない", [
    ("レビューのやり方が大きく変わってしまうのでは?", "変わらない。逆変換ツールにより、これまで通り状態遷移図・判断表・タイミングチャートを見てレビューできる。"),
    ("全部を一気にYAMLへ置き換えるのか?", "本資料では置き換えの是非と考え方の共有が目的。対象範囲・進め方は別途チームで検討する。"),
])

# ============================================================
# Slide 27: FAQ-3
# ============================================================
faq_slide("よくある不安への回答(3/3)", "運用・移行の実務面について", [
    ("ツールの導入・運用コストがかかるのでは?", "既存のGit環境や汎用エディタで始められ、追加コストは小さく抑えられる想定。詳細は別途整理する。"),
    ("既存のExcel資産はどうなるのか?", "既存資産は参照用として残しつつ、対象を絞って段階的にYAML化する。移行方法は別途検討する。"),
])

# ============================================================
# Slide 28: まとめ
# ============================================================
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
add_textbox(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.6), "まとめ", size=28, color=WHITE, bold=True)
data5 = [
    ["現状の課題", "YAML化による対応"],
    ["① 変更点の見落とし", "Git差分による機械的な変更検出"],
    ["② 影響範囲の見極めが困難", "トレーサビリティマトリクスによる自動追跡"],
    ["③ 小手先修正の繰り返し", "構造化された要求から設計・テストを自動生成し、根本対応の土台を作る"],
]
add_table(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(2.3), data5,
          col_widths=[Inches(4.5), Inches(7.0)], row_h=Inches(0.58), body_size=13.5)
add_textbox(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.0),
            "変えるのは「記述の形式」だけ。判断表・状態遷移・シーケンスという設計の考え方、\nレビューで確認する図や表は、これまでと同じものを使い続けられる。",
            size=17, color=WHITE, bold=True)
add_textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2),
            "まずは対象範囲を絞った試行から始め、効果を確認しながら適用範囲を広げていく。\n次のステップについて、皆さんのご意見をお聞かせください。",
            size=15, color=RGBColor(0xB8,0xD4,0xEA))

# ---------- フッター一括付与(表紙・まとめスライドを除く) ----------
slides_list = list(prs.slides)
last_idx = len(slides_list) - 1
for idx, sl in enumerate(slides_list):
    if idx == 0 or idx == last_idx:
        continue
    add_footer(sl, page_number=idx + 1)

prs.save("要求仕様デジタル化による開発プロセス改革.pptx")
print("SAVED. total slides:", len(prs.slides._sldIdLst))





